from pptx import Presentation


def create_presentation():
    prs = Presentation()

    # Slide 1: Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "HSSN 驱动的注意力资本市场"
    slide.placeholders[1].text = "将激励导向 Sonic SVM 程序\n\n团队 & 日期\n\n来源：Notion 页面"

    # Slide 2: 目录
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "目录"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    items = [
        "背景与挑战",
        "HSSN 注意力机制概述",
        "模型架构",
        "Sonic SVM 程序简介",
        "激励导向机制",
        "实验与评估",
        "商业价值与应用场景",
        "总结与展望",
        "Q&A",
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

    # Slide 3: 背景与挑战
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "背景与挑战"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    content = [
        "传统资本市场面临的问题",  
        "  - 信息不对称导致资源配置效率低",  
        "  - 流动性不足与定价偏差",  
        "机器学习在金融中的应用瓶颈",  
        "  - 现有 SVM 模型速度与准确度难以兼顾",  
        "  - 缺乏对市场行为的动态反馈",
    ]
    for line in content:
        p = body.add_paragraph()
        p.text = line
        p.level = 0 if not line.startswith("  -") else 1

    # Slide 4: HSSN 注意力机制概述
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "HSSN 注意力机制概述"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    content = [
        "HSSN（Hierarchical Self-Supervised Network）核心思想:",
        "  - 自监督学习构建分层表示",
        "  - 引入注意力机制强化关键特征",
        "如何提升资本市场效率:",
        "  - 动态捕捉市场热点",
        "  - 实时调整投资策略",
    ]
    for line in content:
        p = body.add_paragraph()
        p.text = line
        p.level = 0 if not line.startswith("  -") else 1

    # Slide 5: 模型架构
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "模型架构"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    content = [
        "总体架构示意:",
        "  - 数据输入层：多源市场数据",
        "  - 表征学习层：HSSN 自监督模块",
        "  - 注意力加权层：重要性打分",
        "  - 输出层：Sonic SVM 优化",
        "核心模块说明:",
        "    1. 特征抽取",
        "    2. 注意力分配",
        "    3. SVM 分类/回归",
    ]
    for line in content:
        p = body.add_paragraph()
        p.text = line
        if line.startswith("  -"):
            p.level = 1
        elif line.strip().startswith(("1.", "2.", "3.")):
            p.level = 2
        else:
            p.level = 0

    # Slide 6: Sonic SVM 程序简介
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Sonic SVM 程序简介"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    content = [
        "Sonic SVM 特点:",
        "  - 轻量级，高并发",
        "  - 支持在线更新",
        "与传统 SVM 对比:",
        "  - 速度提升 3×",
        "  - 精度不降低或略有提升",
    ]
    for line in content:
        p = body.add_paragraph()
        p.text = line
        p.level = 0 if not line.startswith("  -") else 1

    # Slide 7: 激励导向机制
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "激励导向机制"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    content = [
        "设计目标：引导市场参与者行为",
        "激励结构:",
        "  - 成交量奖励",
        "  - 风险控制激励",
        "如何与注意力模型结合:",
        "  - 注意力得分 → 奖励权重",
    ]
    for line in content:
        p = body.add_paragraph()
        p.text = line
        p.level = 0 if not line.startswith("  -") else 1

    # Slide 8: 实验与评估
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "实验与评估"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    content = [
        "实验设置:",
        "  - 数据来源：真实历史交易数据",
        "  - 指标：收益率（ROI）、夏普比率（Sharpe）、回撤率",
        "结果亮点:",
        "  - ROI 提升 12%",
        "  - 夏普比率提升 25%",
        "  - 最大回撤下降 8%",
    ]
    for line in content:
        p = body.add_paragraph()
        p.text = line
        p.level = 0 if not line.startswith("  -") else 1

    # Slide 9: 商业价值与应用场景
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "商业价值与应用场景"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    content = [
        "潜在客户:",
        "  - 对冲基金",
        "  - 量化交易团队",
        "应用场景:",
        "  - 高频交易决策支持",
        "  - 风控预警系统",
        "收益预测:",
        "  - 年化收益率提升空间 >15%",
    ]
    for line in content:
        p = body.add_paragraph()
        p.text = line
        p.level = 0 if not line.startswith("  -") else 1

    # Slide 10: 总结与展望
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "总结与展望"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    content = [
        "研究贡献:",
        "  - 创新性地将 HSSN 与 Sonic SVM 融合",
        "  - 首次在激励导向下实现注意力驱动的市场优化",
        "下一步工作:",
        "  - 扩展至多市场、多品种",
        "  - 深入用户行为建模",
    ]
    for line in content:
        p = body.add_paragraph()
        p.text = line
        p.level = 0 if not line.startswith("  -") else 1

    # Slide 11: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Q&A"
    slide.placeholders[1].text = "欢迎现场提问！"

    # Save presentation to file
    prs.save("HSSN_Sonic_SVM_Presentation.pptx")


if __name__ == "__main__":
    create_presentation() 